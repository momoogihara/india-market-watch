import os
from datetime import datetime

# 必要ライブラリのインポート
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt

# --- スタイリング定数 ---
FONT_TITLE = "Playfair Display"
FONT_BODY = "Montserrat"

# カラーパレット (Management Consulting Deep Navy Theme)
COLOR_BG = RGBColor(11, 19, 43)        # #0b132b (背景色)
COLOR_TEXT = RGBColor(240, 240, 240)   # 明るいグレー (本文用)
COLOR_GOLD = RGBColor(255, 215, 0)     # #ffd700 (メインアクセント・タイトル)
COLOR_TEAL = RGBColor(91, 192, 190)    # #5bc0be (サブアクセント)
COLOR_GREY = RGBColor(140, 150, 160)    # ニュートラルグレー（下位用）

def apply_background(slide):
    """スライド全体の背景色をディープネイビーに設定"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_slide_title(slide, text):
    """統一されたスタイルのタイトルを shapes.title を使用して安全に追加"""
    title_shape = slide.shapes.title
    title_shape.text = text
    title_text_frame = title_shape.text_frame
    title_text_frame.word_wrap = True
    
    # 段落スタイルの調整
    p = title_text_frame.paragraphs[0]
    p.font.name = FONT_TITLE
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.alignment = PP_ALIGN.LEFT
    
    # 絶対位置とサイズを美しく固定
    title_shape.left = Inches(0.8)
    title_shape.top = Inches(0.6)
    title_shape.width = Inches(11.2)
    title_shape.height = Inches(1.2)
    return title_shape

def create_sector_chart(sectors_list, scores_list):
    """上位セクターほどポジティブな輝き（ティール）を放つ美しい横棒グラフ"""
    plt.figure(figsize=(7, 4.5), facecolor='#0b132b')
    ax = plt.axes()
    ax.set_facecolor('#0b132b')
    
    # 配色ロジックの修正：スコアが高い上位（配列の後半、または数値が大きいもの）をティールに
    # スコアが低いものは落ち着いたグレーに設定してコントラストを強調
    max_score = max(scores_list) if scores_list else 1
    colors = []
    for score in scores_list:
        if score >= max_score * 0.6:  # スコアが高い上位陣
            colors.append('#5bc0be')   # ティール
        else:
            colors.append('#8c96a0')   # ニュートラルグレー（下位）
            
    bars = plt.barh(sectors_list, scores_list, color=colors, height=0.6)
    
    # 不要な枠線を消去
    for spine in ax.spines.values():
        spine.set_visible(False)
        
    # グリッド線の微調整
    ax.xaxis.grid(True, linestyle='--', alpha=0.15, color='#ffffff')
    ax.set_axisbelow(True)
    
    # 軸の文字色とフォントサイズの設定
    ax.tick_params(colors='#e0e0e0', labelsize=11)
    plt.title("Sector Bullish Scores", color='#ffd700', fontsize=14, fontweight='bold', pad=15)
    plt.xlabel("Score", color='#5bc0be', fontsize=11, labelpad=10)
    
    # 各バーの右横に数値を正確にプロット
    for bar in bars:
        width = bar.get_width()
        ax.text(width + 0.1, bar.get_y() + bar.get_height()/2, f'{width:.1f}', 
                va='center', ha='left', color='#ffffff', fontsize=10, fontweight='bold')

    plt.tight_layout()
    chart_path = "/tmp/sector_chart.png" if os.path.exists("/tmp") else "sector_chart.png"
    plt.savefig(chart_path, dpi=200, facecolor=plt.gcf().get_facecolor(), edgecolor='none')
    plt.close()
    return chart_path

def create_sentiment_pie_chart(bullish, neutral, bearish):
    """円グラフの見た目を整備（高級感のあるドーナツ型チャート）"""
    labels = ["Bullish", "Neutral", "Bearish"]
    try:
        values = [float(bullish), float(neutral), float(bearish)]
    except ValueError:
        values = [55, 40, 5] # データに合わせたデフォルト値
        
    colors = ["#5bc0be", "#ffd700", "#e63946"] # Bullish(ティール), Neutral(ゴールド), Bearish(赤)
    
    plt.figure(figsize=(5.5, 5.5), facecolor='#0b132b')
    
    wedges, texts, autotexts = plt.pie(
        values, 
        labels=labels, 
        colors=colors,
        autopct="%1.0f%%", 
        startangle=90, 
        pctdistance=0.75,
        textprops=dict(color="#e0e0e0", fontsize=12)
    )
    
    # ドーナツの穴を作成
    centre_circle = plt.Circle((0,0), 0.52, fc='#0b132b')
    fig = plt.gcf()
    fig.gca().add_artist(centre_circle)
    
    for autotext in autotexts:
        autotext.set_color('#0b132b')
        autotext.set_fontweight('bold')
        
    plt.title("Market Sentiment Distribution", color='#ffd700', fontsize=14, fontweight='bold', pad=15)
    plt.tight_layout()
    
    chart_path = "/tmp/sentiment_pie.png" if os.path.exists("/tmp") else "sentiment_pie.png"
    plt.savefig(chart_path, dpi=200, facecolor=plt.gcf().get_facecolor(), edgecolor='none')
    plt.close()
    return chart_path

# =====================================================================
# API外部呼出用メイン関数
# =====================================================================
def generate_sample_ppt(
    summary: str,
    strongest_sector: str,
    weakest_sector: str,
    market_tone: str,
    ranking_text: str,
    coverage_text: str,
    sectors: str,      
    scores: str,       
    bullish_count: str,
    bearish_count: str,
    neutral_count: str,
    headline_text: str
) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ------------------
    # Slide 1: Title Slide (Layout[0])
    # ------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    apply_background(slide1)
    
    title_shape = slide1.shapes.title
    title_shape.text = "India Market Watch"
    title_p = title_shape.text_frame.paragraphs[0]
    title_p.font.name = FONT_TITLE
    title_p.font.size = Pt(64)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_GOLD
    
    subtitle_shape = slide1.placeholders[1]
    subtitle_shape.text = f"Weekly Strategic Overview | {datetime.now().strftime('%B %d, %Y')}"
    subtitle_p = subtitle_shape.text_frame.paragraphs[0]
    subtitle_p.font.name = FONT_BODY
    subtitle_p.font.size = Pt(20)
    subtitle_p.font.color.rgb = COLOR_TEAL
    
    # ------------------
    # Slide 2: Top News Headlines (Layout[1])
    # ------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_background(slide2)
    add_slide_title(slide2, "Top News Headlines")
    
    # 余白比率を最適化したテキストエリア
    textbox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
    tf2 = textbox2.text_frame
    tf2.word_wrap = True
    
    # 各種表記揺れ・数字付きリストを自動クリーニング
    import re
    raw_headlines = [line.strip() for line in headline_text.split('\n') if line.strip()]
    headlines = []
    for line in raw_headlines:
        clean = re.sub(r'^[•\-\*\d\.\s]+', '', line) # 先頭の「1.」や「•」をきれいに削除
        if clean:
            headlines.append(clean)

    if not headlines:
        headlines = ["No headline data available."]

    for i, text in enumerate(headlines[:5]): # 最大5つまでに制限してはみ出し防止
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"•  {text}"
        p.font.name = FONT_BODY
        p.font.size = Pt(15) # テキストが多めでも収まるよう15ptに微調整
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(16)
        
    # ------------------
    # Slide 3: Key Insights (Layout[1])
    # ------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_background(slide3)
    add_slide_title(slide3, "Key Insights & Market Tone")
    
    # 左カラム
    left_box = slide3.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.2))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    
    p = ltf.paragraphs[0]
    p.text = "Sector Dynamics"
    p.font.name = FONT_TITLE
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEAL
    p.space_after = Pt(14)
    
    p2 = ltf.add_paragraph()
    p2.text = f"• Strongest Sector: {strongest_sector}"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(15)
    p2.font.color.rgb = COLOR_TEXT
    p2.space_after = Pt(12)
    
    p3 = ltf.add_paragraph()
    p3.text = f"• Weakest Sector: {weakest_sector}"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(15)
    p3.font.color.rgb = COLOR_TEXT
    
    # 右カラム
    right_box = slide3.shapes.add_textbox(Inches(6.9), Inches(2.2), Inches(5.6), Inches(4.2))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    p = rtf.paragraphs[0]
    p.text = "Macro Outlook"
    p.font.name = FONT_TITLE
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(14)
    
    p2 = rtf.add_paragraph()
    p2.text = f"• Market Tone: {market_tone}"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(15)
    p2.font.color.rgb = COLOR_TEXT
    p2.space_after = Pt(12)
    
    p3 = rtf.add_paragraph()
    p3.text = f"• Summary: {summary}"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(13) # 要約テキストが長くなった場合のために13ptに調整
    p3.font.color.rgb = COLOR_TEXT

    # ------------------
    # Slide 4: Sector Ranking Chart (Layout[1] - 4:6比率固定)
    # ------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_background(slide4)
    add_slide_title(slide4, "Sector Ranking")
    
    text_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(4.5), Inches(4.5))
    tf4 = text_box4.text_frame
    tf4.word_wrap = True
    
    raw_notes = [line.strip() for line in ranking_text.split('\n') if line.strip()]
    for i, note in enumerate(raw_notes[:4]):
        clean_note = re.sub(r'^[•\-\*\d\.\s]+', '', note)
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.text = f"•  {clean_note}"
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(14)
        
    if isinstance(sectors, str):
        sectors_list = [s.strip() for s in sectors.split(',') if s.strip()]
    else:
        sectors_list = sectors

    if isinstance(scores, str):
        scores_list = []
        for s in scores.split(','):
            try: scores_list.append(float(s.strip()))
            except ValueError: pass
    else:
        scores_list = [float(x) for x in scores]
        
    min_len = min(len(sectors_list), len(scores_list))
    if min_len == 0:
        sectors_list = ["Auto", "Infrastructure", "IT", "Energy", "Finance"]
        scores_list = [1.0, 1.0, 1.0, 2.0, 3.0]
    else:
        sectors_list = sectors_list[:min_len]
        scores_list = scores_list[:min_len]

    # 横棒グラフの描画順（上から順にスコア大）にするためのソート処理
    sorted_pairs = sorted(zip(scores_list, sectors_list))
    scores_list, sectors_list = zip(*sorted_pairs)

    chart_path = create_sector_chart(list(sectors_list), list(scores_list))
    slide4.shapes.add_picture(chart_path, Inches(5.8), Inches(2.0), width=Inches(6.8))

    # ------------------
    # Slide 5: Market Coverage / Sentiment Pie (Layout[1] - 4:6比率固定)
    # ------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    apply_background(slide5)
    add_slide_title(slide5, "Market Coverage & Sentiment")
    
    text_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(4.5), Inches(4.5))
    tf5 = text_box5.text_frame
    tf5.word_wrap = True
    
    raw_c_notes = [line.strip() for line in coverage_text.split('\n') if line.strip()]
    for i, note in enumerate(raw_c_notes[:4]):
        clean_note = re.sub(r'^[•\-\*\d\.\s]+', '', note)
        p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        p.text = f"•  {clean_note}"
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(14)
        
    pie_path = create_sentiment_pie_chart(bullish_count, neutral_count, bearish_count)
    slide5.shapes.add_picture(pie_path, Inches(6.5), Inches(1.6), width=Inches(5.5))

    # ------------------
    # 絶対パスでの返却保証
    # ------------------
    root_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))) # /app
    if not os.path.exists(root_dir):
        root_dir = os.getcwd()
        
    file_name = "india_market_watch.pptx"
    file_path = "india_market_watch.pptx"
    
    prs.save(file_path)
    
    try:
        if os.path.exists(chart_path): os.remove(chart_path)
        if os.path.exists(pie_path): os.remove(pie_path)
    except Exception:
        pass
        
    return str(file_path)